"""One-off generator: PORTAL_SLIDE_DECK.pptx from structured content. Run: py -3 build_portal_slide_deck_pptx.py"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUT = Path(__file__).resolve().parent / "PORTAL_SLIDE_DECK.pptx"


def add_title_slide(prs, title: str, subtitle: str):
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title: str, bullets: list[str], notes: str | None = None):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Infinity Nakshatra Society",
        "Resident & Maintenance Portal\n\nPurpose-built web application for housing-society operations: "
        "maintenance collection, owner self-service, and administration—aligned with your master plot register.",
    )

    add_bullet_slide(
        prs,
        "What this platform delivers",
        [
            "Single sign-in for administrators and plot owners",
            "Live operational picture of plots, contacts, maintenance receipts, and dues",
            "Controlled workflows: owners propose payments; admins confirm and post",
            "Optional cloud sync so data and settings are shared across devices",
            "Mobile-friendly layout for day-to-day use",
        ],
    )

    add_bullet_slide(
        prs,
        "Service & commercial model",
        [
            "Hosting: included at no charge on standard static hosting suitable for this portal",
            "Licence use: lifetime use of the delivered application for the scope in this deck",
            "Premium tiers: none — no paid modules or feature paywalls in this programme",
            "First-year enhancement service included: refinements and process alignment (agreed scope)",
            "Legal terms: defined in your project agreement; this deck is a functional summary only",
        ],
    )

    add_bullet_slide(
        prs,
        "Technical foundation (high level)",
        [
            "Master data in a spreadsheet-backed register; portal reads on load and periodic refresh",
            "Optional backend (Apps Script) writes collection logs, shared state, and selected field updates",
            "Onboarding uses your workbook and deployment configuration — no sample personal data required",
        ],
    )

    add_bullet_slide(
        prs,
        "Module map — overview",
        [
            "Admin: dashboard metrics, maintenance analytics, directory, payments, approvals, communications, "
            "documents, banking, reporting",
            "Owner: secure login, payment summary and requests, notices, documents, contacts, banking, optional inbox",
        ],
    )

    add_bullet_slide(
        prs,
        "Authentication & security",
        [
            "Role-based login: Administrator vs Owner",
            "Password management and admin-assisted owner reset",
            "Session handling and secure logout; optional per-owner access restrictions",
            "Optional portal access rules tied to maintenance payment status, with explicit overrides",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — top bar & quick actions",
        [
            "Open the linked master workbook",
            "Add payment (single) and add multiple payments (batch by month)",
            "Download report (maintenance received — structured export)",
            "Profile: admin notes, reset owner password, change password, logout",
            "Optional community link when configured",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — dashboard statistics",
        [
            "Total plots; sold vs unsold",
            "Missing primary contact",
            "Registered owners using the portal",
            "Figures derive from the loaded register and portal usage",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — owner payment approvals",
        [
            "Owners cannot post directly to the official ledger from their login",
            "They submit a maintenance request with month and amounts",
            "Admins review, approve, or reject in a dedicated queue",
            "Months already recorded by admin are locked for duplicate owner submissions",
            "Outcome: audit-friendly flow and fewer data conflicts",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — maintenance collection & filters",
        [
            "Filters: financial year scope (multi-select), owner portal FY, month-in-scope (defaults to all months)",
            "KPIs for selection: received, expected, pending/due, late fee",
            "Shortcuts: payment-completed and payment-pending owner lists",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — payments & ledger alignment",
        [
            "See payment per plot: month-wise status; total received, late fee, counts",
            "Add payment and bulk payment flows; backup to sheet when backend is configured",
            "Plot directory: total amount paid aligned with See payment (admin FY scope)",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — plot directory & access",
        [
            "Search and filters: sold/unsold, contact present/absent, single- vs multi-plot owners",
            "Access column: allow/deny owner login with defaults and overrides",
            "Per-row add payment and see payment; multi-plot condensed view",
            "Unsold plots: add/update rows and sync to workbook when API enabled",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — banking & UPI",
        [
            "Society banking / UPI details maintained centrally",
            "View and edit (admin); owners see Account & UPI with clarity-focused layout",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — notices",
        [
            "Create and publish notices (title, body, audience, optional attachment)",
            "Owners see an up-to-date notices list",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — project documents",
        [
            "Curate project or society documents (name, details, link)",
            "Owner read-only list; optional engagement where enabled",
        ],
    )

    add_bullet_slide(
        prs,
        "Admin — service contacts",
        [
            "Maintain vendors and service numbers",
            "Same contacts visible to owners for self-service coordination",
        ],
    )

    add_bullet_slide(
        prs,
        "Owner portal — home & profile",
        [
            "Welcome context and optional guided tips",
            "Account & UPI; add payment (request / approval flow)",
            "Profile: change password, logout",
            "Messages inbox when the office posts items: read state, mark read",
        ],
    )

    add_bullet_slide(
        prs,
        "Owner portal — notices & information",
        [
            "Notices: date, title, details, attachment",
            "Project documents and service contacts",
        ],
    )

    add_bullet_slide(
        prs,
        "Owner portal — payments summary & history",
        [
            "Summary KPIs for active owner financial year: plots, paid, expected, pending, late fee",
            "Footnotes for due dates and multi-plot totals",
            "Interactive payment history chart; plot picker when one login maps to multiple plots",
        ],
    )

    add_bullet_slide(
        prs,
        "Owner portal — plot detail & updates",
        [
            "Read-only plot and owner fields from the register",
            "Update primary/alternate contact where allowed; optional sync to workbook",
        ],
    )

    add_bullet_slide(
        prs,
        "Cross-cutting services",
        [
            "Automatic refresh of register data (configurable)",
            "Global loading indicator; responsive layout; accessibility-minded control labels",
        ],
    )

    add_bullet_slide(
        prs,
        "First-year enhancement programme (included)",
        [
            "Labels, defaults, and filter behaviour",
            "Minor workflow and copy tweaks",
            "FY/month alignment with society resolutions",
            "Light UX polish after committee feedback",
            "Major new modules or integrations: usually scoped separately",
        ],
    )

    add_bullet_slide(
        prs,
        "Lifetime use — what “no premium” means",
        [
            "Core features in this deck: no subscription paywall for this deployment",
            "Hosting on a suitable static tier: no additional hosting fee from this programme",
            "Domain, spreadsheet quotas, and future capacity: society/host responsibility per your setup",
        ],
    )

    add_bullet_slide(
        prs,
        "Summary",
        [
            "Committee / admin: faster collections visibility, approvals, notices, documents, contacts",
            "Owners: clarity on dues, self-service requests, society information in one place",
            "Society: no premium tier for core features; first-year enhancement; lifetime use of delivered scope",
        ],
    )

    add_title_slide(
        prs,
        "Thank you",
        "Infinity Nakshatra Society — Resident & Maintenance Portal\n\n"
        "Operational clarity • Owner convenience • Sustainable, transparent maintenance management\n\n"
        "This deck describes product capabilities only. No personal data, plot identifiers, screenshots, or source code.",
    )

    prs.save(OUT)
    print(f"Wrote: {OUT}")


if __name__ == "__main__":
    main()
